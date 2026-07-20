"""Document parsers for supported upload formats.

Each parser turns raw file bytes into plain text ready for chunking. PDF/DOCX/
PPTX use their respective libraries (imported lazily so the default TXT/CSV path
and the test-suite need nothing extra); TXT/CSV use only the standard library.

All parsers return a single plain-text string. Page boundaries are preserved for
PDFs by inserting form-feed markers that the chunker maps back to page numbers.
"""

from __future__ import annotations

import csv
import io

# Marker inserted between PDF/PPTX pages so downstream chunking can recover the
# page (or slide) number for citations.
PAGE_BREAK = "\f"

SUPPORTED_FORMATS = ("pdf", "docx", "pptx", "txt", "csv")


def parse_pdf(data: bytes) -> str:
    """Extract text from a PDF, one page per ``PAGE_BREAK`` section."""
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages = [(page.extract_text() or "").strip() for page in reader.pages]
    return PAGE_BREAK.join(pages)


def parse_docx(data: bytes) -> str:
    """Extract text from a Word document (paragraphs + table cells)."""
    from docx import Document as DocxDocument

    document = DocxDocument(io.BytesIO(data))
    parts = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def parse_pptx(data: bytes) -> str:
    """Extract text from a PowerPoint deck, one slide per ``PAGE_BREAK``."""
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    slides = []
    for slide in presentation.slides:
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if text:
                        lines.append(text)
        slides.append("\n".join(lines))
    return PAGE_BREAK.join(slides)


def parse_txt(data: bytes) -> str:
    """Decode a plain-text file (UTF-8 with a permissive fallback)."""
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        return data.decode("latin-1", errors="replace")


def parse_csv(data: bytes) -> str:
    """Flatten CSV rows into readable ``column: value`` lines for indexing."""
    text = parse_txt(data)
    reader = csv.reader(io.StringIO(text))
    rows = list(reader)
    if not rows:
        return ""

    header, *body = rows
    header = [h.strip() for h in header]
    lines: list[str] = []
    for row in body:
        pairs = [
            f"{header[i]}: {value.strip()}"
            for i, value in enumerate(row)
            if i < len(header) and value.strip()
        ]
        if pairs:
            lines.append("; ".join(pairs))
    # If there was no usable body (e.g. header only), fall back to raw text.
    return "\n".join(lines) if lines else text


PARSERS = {
    "pdf": parse_pdf,
    "docx": parse_docx,
    "pptx": parse_pptx,
    "txt": parse_txt,
    "csv": parse_csv,
}


def parse(filename: str, data: bytes) -> str:
    """Route a file to the correct parser based on its extension."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    parser = PARSERS.get(ext)
    if parser is None:
        raise ValueError(
            f"Unsupported format: {ext or filename!r}. "
            f"Supported: {', '.join(SUPPORTED_FORMATS)}."
        )
    return parser(data)
