"""Generators for real, byte-valid fixture files in every supported format.

The point of this module is that the format tests are not allowed to cheat. Each
helper produces genuine file bytes that the corresponding third-party library
would produce, so the parsers are exercised against real containers rather than
a mocked-out stand-in.

DOCX and PPTX are built with the same libraries the parsers read them back with.
The PDF is assembled by hand: no PDF *writer* is in the dependency list, and
adding one purely for tests would be a heavier commitment than the roughly forty
lines below. It is a minimal but standards-valid PDF with a correct cross
reference table, one page per string, and Helvetica text drawn with ``Tj``.
"""

from __future__ import annotations

import csv
import io


def make_txt(text: str) -> bytes:
    """Plain UTF-8 text bytes."""
    return text.encode("utf-8")


def make_csv(header: list[str], rows: list[list[str]]) -> bytes:
    """A real CSV, written by the stdlib csv writer (quoting included)."""
    buffer = io.StringIO()
    writer = csv.writer(buffer)
    writer.writerow(header)
    writer.writerows(rows)
    return buffer.getvalue().encode("utf-8")


def make_docx(paragraphs: list[str], table: list[list[str]] | None = None) -> bytes:
    """A real .docx built with python-docx, optionally containing a table."""
    from docx import Document as DocxDocument

    document = DocxDocument()
    for paragraph in paragraphs:
        document.add_paragraph(paragraph)
    if table:
        docx_table = document.add_table(rows=len(table), cols=len(table[0]))
        for row_index, row in enumerate(table):
            for col_index, value in enumerate(row):
                docx_table.cell(row_index, col_index).text = value

    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def make_pptx(slides: list[tuple[str, str]]) -> bytes:
    """A real .pptx built with python-pptx: one (title, body) pair per slide."""
    from pptx import Presentation

    presentation = Presentation()
    layout = presentation.slide_layouts[1]  # Title and Content
    for title, body in slides:
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = body

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _pdf_escape(text: str) -> str:
    """Escape the three characters that are special inside a PDF string."""
    return text.replace("\\", r"\\").replace("(", r"\(").replace(")", r"\)")


def make_pdf(pages: list[str]) -> bytes:
    """A minimal, valid PDF with one text-bearing page per entry in ``pages``.

    Object layout: 1 = catalog, 2 = page tree, 3 = font, then a (page, content
    stream) pair per page. Byte offsets for the xref table are recorded as the
    objects are appended, because an xref that disagrees with the actual offsets
    is the fastest way to make a reader fall back to scanning or give up.
    """
    objects: list[bytes] = []

    def add(body: bytes) -> int:
        """Append an object body and return its 1-based object number."""
        objects.append(body)
        return len(objects)

    page_object_numbers = []
    # Reserve 1 (catalog) and 2 (pages) so the page objects can point at them.
    add(b"")  # 1: catalog, filled in below
    add(b"")  # 2: page tree, filled in below
    font_number = add(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")

    for text in pages:
        stream = (
            f"BT /F1 12 Tf 72 720 Td ({_pdf_escape(text)}) Tj ET".encode("latin-1")
        )
        content_number = add(
            b"<< /Length "
            + str(len(stream)).encode()
            + b" >>\nstream\n"
            + stream
            + b"\nendstream"
        )
        page_number = add(
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 "
            + str(font_number).encode()
            + b" 0 R >> >> /Contents "
            + str(content_number).encode()
            + b" 0 R >>"
        )
        page_object_numbers.append(page_number)

    kids = b" ".join(f"{n} 0 R".encode() for n in page_object_numbers)
    objects[0] = b"<< /Type /Catalog /Pages 2 0 R >>"
    objects[1] = (
        b"<< /Type /Pages /Kids ["
        + kids
        + b"] /Count "
        + str(len(page_object_numbers)).encode()
        + b" >>"
    )

    out = bytearray(b"%PDF-1.4\n")
    offsets: list[int] = []
    for index, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += f"{index} 0 obj\n".encode() + body + b"\nendobj\n"

    xref_offset = len(out)
    out += f"xref\n0 {len(objects) + 1}\n".encode()
    out += b"0000000000 65535 f \n"
    for offset in offsets:
        out += f"{offset:010d} 00000 n \n".encode()
    out += (
        b"trailer\n<< /Size "
        + str(len(objects) + 1).encode()
        + b" /Root 1 0 R >>\nstartxref\n"
        + str(xref_offset).encode()
        + b"\n%%EOF\n"
    )
    return bytes(out)
